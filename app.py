from __future__ import annotations

import io
import math
import re
from collections import defaultdict
from copy import deepcopy
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Pt


APP_DIR = Path(__file__).resolve().parent
DEFAULT_TEMPLATE = APP_DIR / "Modelo_Devolutiva_SST_Fortneer.pptx"

PRIORITY_ORDER = {"ALTA": 0, "MÉDIA": 1, "BAIXA": 2}
RISK_GROUPS = ["FÍSICO", "QUÍMICO", "BIOLÓGICO", "ERGONÔMICO", "ACIDENTE"]
CONTROL_TYPES = [
    "ELIMINAR/SUBSTITUIR",
    "PROTEÇÃO COLETIVA",
    "MEDIDA ADMINISTRATIVA",
    "PROTEÇÃO INDIVIDUAL",
]

# Cores do template Fortneer para a coluna de prioridade do plano de ação.
PRIORITY_COLORS = {
    "ALTA": RGBColor(0xC9, 0x4A, 0x4A),
    "MÉDIA": RGBColor(0xE8, 0xA2, 0x3A),
    "BAIXA": RGBColor(0x49, 0xA6, 0x6F),
}


# -----------------------------------------------------------------------------
# Utilidades de texto
# -----------------------------------------------------------------------------

def clean_line(value: Any) -> str:
    if value is None:
        return ""
    return re.sub(r"\s+", " ", str(value)).strip()


def safe_str(value: Any, default: str = "") -> str:
    if value is None:
        return default
    try:
        if pd.isna(value):
            return default
    except Exception:
        pass
    text = clean_line(value)
    return text if text else default


def unique_preserve(values: list[str]) -> list[str]:
    seen = set()
    result = []
    for value in values:
        normalized = clean_line(value)
        if not normalized:
            continue
        key = normalized.casefold()
        if key not in seen:
            seen.add(key)
            result.append(normalized)
    return result


def truncate(text: str, max_chars: int) -> str:
    text = clean_line(text)
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


# -----------------------------------------------------------------------------
# Leitura e parsing determinístico do relatório PDF
# -----------------------------------------------------------------------------

@st.cache_data(show_spinner=False)
def extract_pdf_text(pdf_bytes: bytes) -> str:
    """Extrai texto do relatório sem OCR e sem IA."""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages = [page.get_text("text") or "" for page in doc]
    doc.close()

    # O modelo de relatório usado possui sumário na primeira página. Ele é
    # descartado para que SETOR/CARGO do índice não sejam tratados como dados.
    if pages:
        first = pages[0]
        if "Relatório de Empresas Geral" in first and first.count("CARGO:") >= 5:
            pages = pages[1:]

    return "\n".join(pages)


def first_match(text: str, pattern: str) -> str:
    match = re.search(pattern, text, flags=re.I | re.M)
    return clean_line(match.group(1)) if match else ""


def parse_company(text: str) -> dict[str, str]:
    return {
        "empresa": first_match(text, r"^EMPRESA:\s*(.+)$"),
        "cnpj": first_match(text, r"^CNPJ:\s*(.+)$"),
        "num_empregados": first_match(text, r"^N[º°]\s*Empregados:\s*(.+)$"),
        "cnae": first_match(text, r"^CNAE:\s*(.+)$"),
        "atividade_principal": first_match(text, r"^Atividade Principal:\s*(.+)$"),
        "grau_risco": first_match(text, r"^Grau de Risco:\s*(.+)$"),
    }


def split_cargo_blocks(text: str) -> list[tuple[str, str, str]]:
    """Separa o relatório em blocos (setor, cargo, conteúdo do cargo)."""
    current_sector = ""
    current_cargo: str | None = None
    current_lines: list[str] = []
    result: list[tuple[str, str, str]] = []

    for raw_line in text.splitlines():
        line = clean_line(raw_line)
        if not line:
            if current_cargo:
                current_lines.append("")
            continue

        sector_match = re.match(r"^SETOR:\s*(.+)$", line, flags=re.I)
        if sector_match:
            if current_cargo is not None:
                result.append((current_sector, current_cargo, "\n".join(current_lines)))
                current_cargo = None
                current_lines = []
            current_sector = clean_line(sector_match.group(1))
            continue

        # Aceita tanto "CARGO: ..." quanto "1.1.1 CARGO: ...".
        cargo_match = re.match(
            r"^(?:\d+(?:\.\d+){2}\s*)?CARGO:\s*(.+)$", line, flags=re.I
        )
        if cargo_match:
            if current_cargo is not None:
                result.append((current_sector, current_cargo, "\n".join(current_lines)))
            current_cargo = clean_line(cargo_match.group(1))
            current_lines = []
            continue

        if current_cargo is not None:
            current_lines.append(line)

    if current_cargo is not None:
        result.append((current_sector, current_cargo, "\n".join(current_lines)))

    return result


def parse_agent_segments(block: str) -> list[dict[str, Any]]:
    """
    Divide cada cargo pelos marcadores AGENTE e captura, para cada agente,
    Conclusão do Laudo e CONCLUSÃO LTCAT/GFIP.
    """
    matches = list(re.finditer(r"(?im)^AGENTE:\s*(.+)$", block))
    segments: list[dict[str, Any]] = []

    for index, match in enumerate(matches):
        agent_name = clean_line(match.group(1))
        end = matches[index + 1].start() if index + 1 < len(matches) else len(block)
        segment = block[match.end() : end]

        conclusion_match = re.search(
            r"(?is)Conclusão do Laudo:\s*(.*?)"
            r"(?=\nCONCLUSÃO LTCAT:|\nRiscos |\nAtestados \(ASO\)|"
            r"\nExames Obrigatórios|\nEPI'?s Recomendados|$)",
            segment,
        )
        conclusion = clean_line(conclusion_match.group(1)) if conclusion_match else ""

        ltcat_match = re.search(
            r"(?is)CONCLUSÃO LTCAT:\s*(.*?)"
            r"(?=\n(?:AGENTE:|Riscos |Atestados \(ASO\)|Exames Obrigatórios|"
            r"EPI'?s Recomendados)|$)",
            segment,
        )
        ltcat = clean_line(ltcat_match.group(1)) if ltcat_match else ""

        gfip_codes = re.findall(r"(?i)GFIP\s*=\s*([0-9])", segment)
        if not gfip_codes:
            # Tolerância a pequenas quebras/ruídos de extração do PDF.
            gfip_codes = re.findall(r"(?i)GFIP\s*[^0-9]{0,8}([0-9])", segment)

        segments.append(
            {
                "agente": agent_name,
                "conclusao": conclusion,
                "ltcat": ltcat,
                "gfips": gfip_codes,
            }
        )

    return segments


def parse_insalubridade(agent_segments: list[dict[str, Any]]) -> tuple[str, str]:
    positives: list[tuple[str, str]] = []
    negative_found = False
    raw_conclusions: list[str] = []

    for item in agent_segments:
        conclusion = safe_str(item.get("conclusao"))
        agent = safe_str(item.get("agente"), "Agente não identificado")
        if not conclusion:
            continue

        raw_conclusions.append(f"{agent}: {conclusion}")
        upper = conclusion.upper()

        if "NÃO INSALUB" in upper or "NAO INSALUB" in upper:
            negative_found = True
            continue

        degree_match = re.search(
            r"INSALUBRES?\s+EM\s+GRAU\s+(MÍNIMO|MINIMO|MÉDIO|MEDIO|MÁXIMO|MAXIMO)",
            upper,
        )
        if degree_match:
            degree = degree_match.group(1)
            degree = {
                "MINIMO": "MÍNIMO",
                "MEDIO": "MÉDIO",
                "MAXIMO": "MÁXIMO",
            }.get(degree, degree)
            positives.append((degree, agent))
        elif "INSALUBR" in upper:
            positives.append(("SIM", agent))

    if positives:
        parts = []
        seen = set()
        for degree, agent in positives:
            key = (degree, agent.casefold())
            if key in seen:
                continue
            seen.add(key)
            if degree == "SIM":
                parts.append(f"Sim — {agent}")
            else:
                parts.append(f"Grau {degree.title()} — {agent}")
        summary = "; ".join(parts)
    elif negative_found:
        summary = "Não insalubre"
    else:
        summary = "Não identificado"

    return summary, "\n".join(raw_conclusions)


def parse_gfip(agent_segments: list[dict[str, Any]]) -> tuple[str, str]:
    pairs: list[tuple[str, str]] = []
    for item in agent_segments:
        agent = safe_str(item.get("agente"), "Agente não identificado")
        for code in item.get("gfips", []):
            pairs.append((str(code), agent))

    if not pairs:
        return "Não identificado", ""

    codes = sorted({code for code, _ in pairs}, key=lambda x: int(x))
    summary = ", ".join(f"GFIP {code}" for code in codes)
    detail = "; ".join(f"GFIP {code} — {agent}" for code, agent in pairs)
    return summary, detail


def parse_exams(block: str) -> list[dict[str, str]]:
    """Captura ASO e Exames Obrigatórios, mantendo a periodicidade associada."""
    lines = [clean_line(line) for line in block.splitlines() if clean_line(line)]
    exams: list[dict[str, str]] = []
    current: dict[str, str] | None = None

    for line in lines:
        aso_match = re.match(r"(?i)^Atestado \(ASO\):\s*(.+)$", line)
        exam_match = re.match(r"(?i)^Exame:\s*(.+)$", line)
        periodicity_match = re.match(r"(?i)^Periodicidade:\s*(.+)$", line)

        if aso_match:
            current = {"exame": clean_line(aso_match.group(1)), "periodicidade": ""}
            exams.append(current)
        elif exam_match:
            current = {"exame": clean_line(exam_match.group(1)), "periodicidade": ""}
            exams.append(current)
        elif periodicity_match and current is not None and not current["periodicidade"]:
            current["periodicidade"] = clean_line(periodicity_match.group(1))

    # Remove duplicatas exatas mantendo a ordem.
    result: list[dict[str, str]] = []
    seen = set()
    for item in exams:
        key = (item["exame"].casefold(), item["periodicidade"].casefold())
        if key not in seen:
            seen.add(key)
            result.append(item)
    return result


def parse_report(text: str) -> tuple[dict[str, str], pd.DataFrame, pd.DataFrame]:
    company = parse_company(text)
    framing_rows: list[dict[str, str]] = []
    exam_rows: list[dict[str, str]] = []

    for sector, role, block in split_cargo_blocks(text):
        agents = parse_agent_segments(block)
        insalubridade, insal_detail = parse_insalubridade(agents)
        gfip, gfip_detail = parse_gfip(agents)

        framing_rows.append(
            {
                "Setor": sector,
                "Cargo": role,
                "Insalubridade": insalubridade,
                "GFIP/LTCAT": gfip,
                "Detalhe insalubridade": insal_detail,
                "Detalhe GFIP/LTCAT": gfip_detail,
            }
        )

        for exam in parse_exams(block):
            exam_rows.append(
                {
                    "Setor": sector,
                    "Cargo": role,
                    "Exame": exam["exame"],
                    "Periodicidade": exam["periodicidade"],
                }
            )

    return company, pd.DataFrame(framing_rows), pd.DataFrame(exam_rows)


# -----------------------------------------------------------------------------
# Formatação dos dados para a apresentação
# -----------------------------------------------------------------------------

def abbreviate_periodicity(text: str) -> str:
    """Abrevia apenas para caber nos cards da apresentação; a tabela segue completa."""
    value = safe_str(text)
    replacements = [
        (r"ADMISSIONAL", "ADM"),
        (r"ANUAL", "AN"),
        (r"MUDANÇA DE FUNÇÃO", "MF"),
        (r"MUDANCA DE FUNCAO", "MF"),
        (r"RETORNO AO TRABALHO", "RT"),
        (r"DEMISSIONAL", "DEM"),
        (r"BIENAL", "BIENAL"),
        (r"SEMESTRAL", "SEM"),
    ]
    upper = value.upper()
    for pattern, replacement in replacements:
        upper = re.sub(pattern, replacement, upper)
    upper = re.sub(r"\s*/\s*", "/", upper)
    return upper


def build_exam_index(exams_df: pd.DataFrame) -> dict[tuple[str, str], list[dict[str, str]]]:
    index: dict[tuple[str, str], list[dict[str, str]]] = defaultdict(list)
    if exams_df is None or exams_df.empty:
        return index

    for _, row in exams_df.iterrows():
        sector = safe_str(row.get("Setor"))
        role = safe_str(row.get("Cargo"))
        exam = safe_str(row.get("Exame"))
        periodicity = safe_str(row.get("Periodicidade"))
        if sector and role and exam:
            index[(sector, role)].append({"exame": exam, "periodicidade": periodicity})
    return index


def compact_exam_count(exams: list[dict[str, str]]) -> str:
    names = unique_preserve([item["exame"] for item in exams])
    if not names:
        return "—"
    if len(names) == 1:
        return truncate(names[0], 28)
    complementaries = max(0, len(names) - 1)
    first = names[0]
    return truncate(f"{first} + {complementaries} compl.", 30)


def format_exam_groups(exams: list[dict[str, str]]) -> str:
    if not exams:
        return "Nenhum exame identificado no relatório."

    grouped: dict[str, list[str]] = defaultdict(list)
    for item in exams:
        periodicity = abbreviate_periodicity(item.get("periodicidade", "")) or "Periodicidade não informada"
        grouped[periodicity].append(safe_str(item.get("exame")))

    lines = []
    for periodicity, names in grouped.items():
        names = unique_preserve(names)
        line = f"{periodicity}: {', '.join(names)}"
        lines.append(line)

    text = "\n".join(lines)
    return truncate(text, 420)


def normalize_actions(actions_df: pd.DataFrame) -> pd.DataFrame:
    columns = [
        "Grupo de risco",
        "Risco/condição",
        "Setor/Cargo afetado",
        "Ação",
        "Prioridade",
        "Tipo de controle",
        "Responsável",
        "Prazo (dias)",
        "Status",
        "Evidência esperada",
    ]

    if actions_df is None or actions_df.empty:
        return pd.DataFrame(columns=columns)

    df = actions_df.copy()
    for column in columns:
        if column not in df.columns:
            df[column] = ""

    df = df[columns]
    df = df[df["Ação"].apply(lambda value: bool(safe_str(value)))].copy()

    def normalize_priority(value: Any) -> str:
        text = safe_str(value, "MÉDIA").upper()
        return text if text in PRIORITY_ORDER else "MÉDIA"

    df["Prioridade"] = df["Prioridade"].apply(normalize_priority)
    df["Status"] = df["Status"].apply(lambda value: safe_str(value, "Aberta"))
    df["Tipo de controle"] = df["Tipo de controle"].apply(safe_str)
    df["Grupo de risco"] = df["Grupo de risco"].apply(lambda value: safe_str(value).upper())

    # Prazo numérico; vazio permanece NaN.
    df["Prazo (dias)"] = pd.to_numeric(df["Prazo (dias)"], errors="coerce")

    df["_priority_order"] = df["Prioridade"].map(PRIORITY_ORDER).fillna(9)
    df = df.sort_values(by=["_priority_order", "Prazo (dias)"], na_position="last")
    return df.drop(columns=["_priority_order"])


# -----------------------------------------------------------------------------
# Manipulação do template PPTX
# -----------------------------------------------------------------------------

def shapes_named(slide, name: str):
    return [shape for shape in slide.shapes if shape.name == name]


def first_named(slide, name: str):
    found = shapes_named(slide, name)
    return found[0] if found else None


def capture_style(shape):
    tf = shape.text_frame
    base_paragraph = tf.paragraphs[0] if tf.paragraphs else None
    base_run = None
    if base_paragraph:
        for run in base_paragraph.runs:
            if run.text or base_run is None:
                base_run = run
                if run.text:
                    break

    style = {
        "alignment": getattr(base_paragraph, "alignment", None) if base_paragraph else None,
        "level": getattr(base_paragraph, "level", 0) if base_paragraph else 0,
        "space_before": getattr(base_paragraph, "space_before", None) if base_paragraph else None,
        "space_after": getattr(base_paragraph, "space_after", None) if base_paragraph else None,
        "line_spacing": getattr(base_paragraph, "line_spacing", None) if base_paragraph else None,
        "font_name": None,
        "font_size": None,
        "bold": None,
        "italic": None,
        "underline": None,
        "color_type": None,
        "rgb": None,
        "theme_color": None,
        "brightness": None,
    }

    if base_run is not None:
        font = base_run.font
        style["font_name"] = font.name
        style["font_size"] = font.size
        style["bold"] = font.bold
        style["italic"] = font.italic
        style["underline"] = font.underline
        try:
            style["color_type"] = font.color.type
            if font.color.type == MSO_COLOR_TYPE.RGB:
                style["rgb"] = font.color.rgb
            elif font.color.type == MSO_COLOR_TYPE.SCHEME:
                style["theme_color"] = font.color.theme_color
                style["brightness"] = font.color.brightness
        except Exception:
            pass

    return style


def apply_run_style(run, style: dict[str, Any], font_size: float | None = None, bold: bool | None = None):
    font = run.font
    if style.get("font_name"):
        font.name = style["font_name"]
    if font_size is not None:
        font.size = Pt(font_size)
    elif style.get("font_size") is not None:
        font.size = style["font_size"]
    if bold is not None:
        font.bold = bold
    elif style.get("bold") is not None:
        font.bold = style["bold"]
    if style.get("italic") is not None:
        font.italic = style["italic"]
    if style.get("underline") is not None:
        font.underline = style["underline"]

    try:
        if style.get("color_type") == MSO_COLOR_TYPE.RGB and style.get("rgb") is not None:
            font.color.rgb = style["rgb"]
        elif style.get("color_type") == MSO_COLOR_TYPE.SCHEME and style.get("theme_color") is not None:
            font.color.theme_color = style["theme_color"]
            if style.get("brightness") is not None:
                font.color.brightness = style["brightness"]
    except Exception:
        pass


def set_shape_text(shape, text: Any, font_size: float | None = None, bold: bool | None = None):
    if shape is None or not shape.has_text_frame:
        return

    value = safe_str(text)
    style = capture_style(shape)
    tf = shape.text_frame
    tf.clear()

    lines = value.split("\n") if value else [""]
    first_paragraph = tf.paragraphs[0]

    for index, line in enumerate(lines):
        paragraph = first_paragraph if index == 0 else tf.add_paragraph()
        try:
            paragraph.alignment = style.get("alignment")
            paragraph.level = style.get("level", 0)
            if style.get("space_before") is not None:
                paragraph.space_before = style["space_before"]
            if style.get("space_after") is not None:
                paragraph.space_after = style["space_after"]
            if style.get("line_spacing") is not None:
                paragraph.line_spacing = style["line_spacing"]
        except Exception:
            pass

        run = paragraph.add_run()
        run.text = line
        apply_run_style(run, style, font_size=font_size, bold=bold)


def set_named_text(slide, name: str, text: Any, font_size: float | None = None, bold: bool | None = None):
    shape = first_named(slide, name)
    if shape is not None:
        set_shape_text(shape, text, font_size=font_size, bold=bold)


def remove_shape(shape):
    element = shape.element
    element.getparent().remove(element)


def add_logo_to_slide(slide, logo_bytes: bytes):
    placeholder = first_named(slide, "logo-placeholder")
    if placeholder is None:
        return

    x, y, w, h = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    try:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), x, y, w, h)
        remove_shape(placeholder)
        logo_text = first_named(slide, "logo-text")
        if logo_text is not None:
            remove_shape(logo_text)
    except Exception:
        # Em caso de imagem inválida, mantém o placeholder em vez de falhar o PPT.
        return


def clone_slide(prs: Presentation, source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Remove eventuais shapes do layout novo e copia todos os shapes do modelo.
    for shape in list(new_slide.shapes):
        remove_shape(shape)
    for shape in source_slide.shapes:
        element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(element, "p:extLst")

    return new_slide


def slide_id_element(prs: Presentation, slide):
    slide_id = prs.part.slide_id(slide.part)
    for element in prs.slides._sldIdLst:
        if element.id == slide_id:
            return element
    raise ValueError("Slide não encontrado na lista da apresentação")


def move_slide_after(prs: Presentation, slide, after_slide):
    lst = prs.slides._sldIdLst
    element = slide_id_element(prs, slide)
    after_element = slide_id_element(prs, after_slide)
    lst.remove(element)
    current = list(lst)
    position = current.index(after_element)
    lst.insert(position + 1, element)


def duplicate_sequence(prs: Presentation, source_slide, total_slides: int) -> list:
    slides = [source_slide]
    previous = source_slide
    for _ in range(max(0, total_slides - 1)):
        clone = clone_slide(prs, source_slide)
        move_slide_after(prs, clone, previous)
        slides.append(clone)
        previous = clone
    return slides


def replace_common_placeholders(prs: Presentation, company: str, period: str):
    for slide in prs.slides:
        # Footers e campos de capa são tratados por nome e também por texto residual.
        if first_named(slide, "footer") is not None:
            set_named_text(slide, "footer", f"{company}  •  {period}")

        if first_named(slide, "company") is not None:
            set_named_text(slide, "company", company)
        if first_named(slide, "date") is not None:
            set_named_text(slide, "date", period)

        # Remove instruções de edição do template que não devem ir ao cliente.
        for note_name in ["matrix-note", "action-note"]:
            shape = first_named(slide, note_name)
            if shape is not None:
                set_shape_text(shape, "")


def update_slide_numbers(prs: Presentation):
    for index, slide in enumerate(prs.slides, start=1):
        number_text = f"{index:02d}"
        if first_named(slide, "slide-number") is not None:
            set_named_text(slide, "slide-number", number_text)
        if first_named(slide, "section-number") is not None:
            set_named_text(slide, "section-number", number_text)


def fill_scope_slide(slide, company_info: dict[str, str], framing_df: pd.DataFrame, settings: dict[str, Any]):
    metrics = shapes_named(slide, "metric-value")
    employees = safe_str(company_info.get("num_empregados"), "0")
    sectors = framing_df["Setor"].nunique() if not framing_df.empty else 0
    roles = len(framing_df) if not framing_df.empty else 0
    units = int(settings.get("units", 1) or 1)
    values = [employees, str(sectors), str(roles), str(units)]
    for shape, value in zip(metrics, values):
        set_shape_text(shape, value)

    bullet_texts = shapes_named(slide, "bullet-text")
    scope_values = [
        f"Documentos considerados: {settings.get('documents', 'Relatório Geral / PGR')}",
        f"Visitas e levantamentos realizados em: {settings.get('visits', 'conforme documentação técnica')}",
        f"Responsáveis técnicos: {settings.get('technical_responsibles', 'conforme documentos vigentes')}",
    ]
    for shape, value in zip(bullet_texts, scope_values):
        set_shape_text(shape, value, font_size=10.5)

    set_named_text(slide, "note", "")


def fill_status_slide(slide, actions_df: pd.DataFrame, settings: dict[str, Any]):
    status_values = shapes_named(slide, "status-value")
    action_count = len(actions_df)
    values = [
        safe_str(settings.get("documentation_status"), "Relatório processado"),
        safe_str(settings.get("exam_status"), "Matriz extraída do relatório"),
        safe_str(settings.get("risk_control_status"), "Ver riscos e plano de ação"),
        f"{action_count} ações cadastradas",
    ]
    for shape, value in zip(status_values, values):
        set_shape_text(shape, value, font_size=11)

    priority_action = "Nenhuma ação cadastrada"
    if not actions_df.empty:
        high = actions_df[actions_df["Prioridade"] == "ALTA"]
        selected = high.iloc[0] if not high.empty else actions_df.iloc[0]
        priority_action = safe_str(selected.get("Ação"), priority_action)

    set_named_text(slide, "priority-title", truncate(priority_action, 130), font_size=17)
    set_named_text(
        slide,
        "priority-copy",
        safe_str(settings.get("priority_impact"), "Impacto esperado: redução da exposição e regularização das pendências priorizadas."),
        font_size=10.5,
    )
    set_named_text(slide, "legend", "")


def fill_risk_slide(slide, actions_df: pd.DataFrame):
    agents = shapes_named(slide, "risk-agent")
    scopes = shapes_named(slide, "risk-scope")

    for group, agent_shape, scope_shape in zip(RISK_GROUPS, agents, scopes):
        subset = actions_df[actions_df["Grupo de risco"] == group] if not actions_df.empty else pd.DataFrame()
        if subset.empty:
            risk_text = "—"
            scope_text = "—"
        else:
            risk_text = "; ".join(unique_preserve([safe_str(v) for v in subset["Risco/condição"].tolist()]))
            scope_text = "; ".join(unique_preserve([safe_str(v) for v in subset["Setor/Cargo afetado"].tolist()]))
            risk_text = truncate(risk_text or "Ações cadastradas", 70)
            scope_text = truncate(scope_text or "Escopo informado no plano", 70)

        set_shape_text(agent_shape, risk_text, font_size=10)
        set_shape_text(scope_shape, scope_text, font_size=9.5)


def fill_function_slides(slides: list, framing_df: pd.DataFrame, exams_df: pd.DataFrame):
    exam_index = build_exam_index(exams_df)
    records = framing_df.to_dict("records") if not framing_df.empty else []

    for page_index, slide in enumerate(slides):
        set_named_text(slide, "title", "Enquadramentos por setor e função")
        headers = shapes_named(slide, "tht")
        header_values = ["FUNÇÃO", "SETOR", "EXAMES", "INSALUBRIDADE", "GFIP / LTCAT"]
        for shape, value in zip(headers, header_values):
            set_shape_text(shape, value, font_size=9.8)

        cells = shapes_named(slide, "celltxt")
        chunk = records[page_index * 5 : (page_index + 1) * 5]

        for row_index in range(5):
            values = ["", "", "", "", ""]
            if row_index < len(chunk):
                record = chunk[row_index]
                key = (safe_str(record.get("Setor")), safe_str(record.get("Cargo")))
                exams = exam_index.get(key, [])
                values = [
                    safe_str(record.get("Cargo"), "—"),
                    safe_str(record.get("Setor"), "—"),
                    compact_exam_count(exams),
                    safe_str(record.get("Insalubridade"), "—"),
                    safe_str(record.get("GFIP/LTCAT"), "—"),
                ]

            for col_index, value in enumerate(values):
                cell = cells[row_index * 5 + col_index]
                font_size = 9.3
                if col_index in (2, 3, 4):
                    font_size = 8.8
                set_shape_text(cell, truncate(value, 58), font_size=font_size)


def build_role_exam_records(framing_df: pd.DataFrame, exams_df: pd.DataFrame) -> list[dict[str, str]]:
    exam_index = build_exam_index(exams_df)
    records = []
    if framing_df.empty:
        return records

    for _, row in framing_df.iterrows():
        sector = safe_str(row.get("Setor"))
        role = safe_str(row.get("Cargo"))
        records.append(
            {
                "Setor": sector,
                "Cargo": role,
                "Exames": format_exam_groups(exam_index.get((sector, role), [])),
            }
        )
    return records


def fill_exam_slides(slides: list, framing_df: pd.DataFrame, exams_df: pd.DataFrame):
    records = build_role_exam_records(framing_df, exams_df)

    for page_index, slide in enumerate(slides):
        roles = shapes_named(slide, "exam-role")
        lists = shapes_named(slide, "exam-list")
        chunk = records[page_index * 3 : (page_index + 1) * 3]

        for row_index in range(3):
            if row_index < len(chunk):
                record = chunk[row_index]
                role_text = f"{record['Cargo']} — {record['Setor']}"
                exam_text = record["Exames"]
            else:
                role_text = ""
                exam_text = ""

            set_shape_text(roles[row_index], truncate(role_text, 90), font_size=10.3)
            set_shape_text(lists[row_index], exam_text, font_size=8.6)

        set_named_text(
            slide,
            "exam-alert-copy",
            "Abreviações: ADM admissional • AN anual • MF mudança de função • RT retorno ao trabalho • DEM demissional.",
            font_size=9.2,
        )


def extract_nontrivial_gfip(gfip_summary: str) -> list[str]:
    codes = re.findall(r"GFIP\s*([0-9])", safe_str(gfip_summary), flags=re.I)
    return sorted({code for code in codes if code != "1"}, key=int)


def fill_legal_summary_slide(slide, framing_df: pd.DataFrame):
    positives = []
    previd = []

    if not framing_df.empty:
        for _, row in framing_df.iterrows():
            sector = safe_str(row.get("Setor"))
            role = safe_str(row.get("Cargo"))
            insal = safe_str(row.get("Insalubridade"))
            gfip = safe_str(row.get("GFIP/LTCAT"))

            if insal and insal not in {"Não insalubre", "Não identificado"}:
                positives.append(f"• {role} ({sector}) — {insal}")

            special_codes = extract_nontrivial_gfip(gfip)
            if special_codes:
                previd.append(f"• {role} ({sector}) — " + ", ".join(f"GFIP {c}" for c in special_codes))

    left_text = "Conclusão do laudo:\n" + ("\n".join(positives[:8]) if positives else "Nenhum enquadramento positivo identificado.")
    if len(positives) > 8:
        left_text += f"\n• +{len(positives) - 8} registros nos slides por função"

    right_text = "Conclusão técnica:\n" + ("\n".join(previd[:8]) if previd else "Nenhum código GFIP diferente de 1 identificado.")
    if len(previd) > 8:
        right_text += f"\n• +{len(previd) - 8} registros nos slides por função"

    set_named_text(slide, "left-result", left_text, font_size=9.7)
    set_named_text(slide, "right-result", right_text, font_size=9.7)


def fill_control_hierarchy_slide(slide, actions_df: pd.DataFrame):
    mapping = {
        "ELIMINAR/SUBSTITUIR": "ELIMINAR OU SUBSTITUIR",
        "PROTEÇÃO COLETIVA": "PROTEÇÃO COLETIVA",
        "MEDIDA ADMINISTRATIVA": "MEDIDA ADMINISTRATIVA",
        "PROTEÇÃO INDIVIDUAL": "PROTEÇÃO INDIVIDUAL",
    }
    heads = shapes_named(slide, "control-head")
    copies = shapes_named(slide, "control-copy")

    for index, control_type in enumerate(CONTROL_TYPES):
        set_shape_text(heads[index], mapping[control_type], font_size=10.3)
        if actions_df.empty:
            text = "Sem ação cadastrada."
        else:
            subset = actions_df[actions_df["Tipo de controle"] == control_type]
            names = unique_preserve([safe_str(v) for v in subset["Ação"].tolist()])
            text = "\n".join(f"• {truncate(name, 70)}" for name in names[:3]) or "Sem ação cadastrada."
            if len(names) > 3:
                text += f"\n• +{len(names) - 3} ações"
        set_shape_text(copies[index], text, font_size=9.2)


def format_deadline(days: Any) -> str:
    try:
        if pd.isna(days):
            return "A definir"
    except Exception:
        pass
    try:
        value = int(float(days))
        return f"{value} dias"
    except Exception:
        return safe_str(days, "A definir")


def fill_action_slides(slides: list, actions_df: pd.DataFrame):
    records = actions_df.to_dict("records") if not actions_df.empty else []

    for page_index, slide in enumerate(slides):
        texts = shapes_named(slide, "acellt")
        backgrounds = shapes_named(slide, "acell")
        chunk = records[page_index * 5 : (page_index + 1) * 5]

        for row_index in range(5):
            if row_index < len(chunk):
                record = chunk[row_index]
                priority = safe_str(record.get("Prioridade"), "MÉDIA").upper()
                values = [
                    priority,
                    safe_str(record.get("Ação"), "—"),
                    safe_str(record.get("Responsável"), "A definir"),
                    format_deadline(record.get("Prazo (dias)")),
                    safe_str(record.get("Status"), "Aberta"),
                ]
            else:
                priority = ""
                values = ["", "", "", "", ""]

            for col_index, value in enumerate(values):
                cell = texts[row_index * 5 + col_index]
                font_size = 9.0 if col_index == 1 else 9.3
                set_shape_text(cell, truncate(value, 84), font_size=font_size)

            # Cada linha possui 5 backgrounds; altera somente o primeiro (NÍVEL).
            priority_background = backgrounds[row_index * 5]
            if priority in PRIORITY_COLORS:
                priority_background.fill.solid()
                priority_background.fill.fore_color.rgb = PRIORITY_COLORS[priority]


def phase_for_days(days: Any) -> str:
    try:
        if pd.isna(days):
            return "PRÓXIMO CICLO"
    except Exception:
        pass

    try:
        value = int(float(days))
    except Exception:
        return "PRÓXIMO CICLO"

    if value <= 30:
        return "0–30 DIAS"
    if value <= 60:
        return "31–60 DIAS"
    if value <= 90:
        return "61–90 DIAS"
    return "PRÓXIMO CICLO"


def fill_timeline_slide(slide, actions_df: pd.DataFrame, settings: dict[str, Any]):
    phase_titles = shapes_named(slide, "phase-title")
    phase_copies = shapes_named(slide, "phase-copy")

    grouped: dict[str, list[dict[str, Any]]] = {title: [] for title in [
        "0–30 DIAS", "31–60 DIAS", "61–90 DIAS", "PRÓXIMO CICLO"
    ]}

    if not actions_df.empty:
        for record in actions_df.to_dict("records"):
            grouped[phase_for_days(record.get("Prazo (dias)"))].append(record)

    for title_shape, copy_shape in zip(phase_titles, phase_copies):
        phase = safe_str(title_shape.text)
        records = grouped.get(phase, [])
        if not records:
            text = "Sem ação cadastrada neste ciclo."
        else:
            lines = []
            for record in records[:5]:
                action = truncate(safe_str(record.get("Ação"), "Ação"), 72)
                lines.append(f"• {action}")
            evidences = unique_preserve([safe_str(record.get("Evidência esperada")) for record in records])
            if evidences:
                lines.append("Evidência: " + truncate("; ".join(evidences[:3]), 120))
            if len(records) > 5:
                lines.append(f"+{len(records) - 5} ações neste ciclo")
            text = "\n".join(lines)
        set_shape_text(copy_shape, text, font_size=9.0)

    review_date = safe_str(settings.get("review_date"), "a definir")
    participants = safe_str(settings.get("review_participants"), "responsáveis definidos no plano")
    set_named_text(
        slide,
        "review-copy",
        f"Revisão de acompanhamento prevista para: {review_date}  •  Participantes: {participants}",
        font_size=9.2,
    )


def fill_contact_slide(slide, settings: dict[str, Any]):
    contact = "\n".join(
        [
            safe_str(settings.get("contact_name"), "Responsável Fortneer"),
            safe_str(settings.get("contact_phone"), "Telefone"),
            safe_str(settings.get("contact_email"), "E-mail"),
        ]
    )
    set_named_text(slide, "contact-copy", contact, font_size=10.5)


def build_presentation(
    template_bytes: bytes,
    company_info: dict[str, str],
    framing_df: pd.DataFrame,
    exams_df: pd.DataFrame,
    actions_df: pd.DataFrame,
    settings: dict[str, Any],
    logo_bytes: bytes | None = None,
) -> bytes:
    prs = Presentation(io.BytesIO(template_bytes))

    if len(prs.slides) < 15:
        raise ValueError("O template precisa manter a estrutura de 15 slides do modelo Fortneer.")

    # Referências aos slides-base antes de qualquer duplicação.
    slide_scope = prs.slides[2]       # 03
    slide_status = prs.slides[3]      # 04
    slide_risks = prs.slides[5]       # 06
    slide_functions = prs.slides[6]   # 07
    slide_exams = prs.slides[8]       # 09
    slide_legal = prs.slides[9]       # 10
    slide_controls = prs.slides[10]   # 11
    slide_actions = prs.slides[12]    # 13
    slide_timeline = prs.slides[13]   # 14
    slide_contact = prs.slides[14]    # 15

    actions = normalize_actions(actions_df)

    function_slide_count = max(1, math.ceil(max(len(framing_df), 1) / 5))
    exam_slide_count = max(1, math.ceil(max(len(framing_df), 1) / 3))
    action_slide_count = max(1, math.ceil(max(len(actions), 1) / 5))

    function_slides = duplicate_sequence(prs, slide_functions, function_slide_count)
    exam_slides = duplicate_sequence(prs, slide_exams, exam_slide_count)
    action_slides = duplicate_sequence(prs, slide_actions, action_slide_count)

    fill_scope_slide(slide_scope, company_info, framing_df, settings)
    fill_status_slide(slide_status, actions, settings)
    fill_risk_slide(slide_risks, actions)
    fill_function_slides(function_slides, framing_df, exams_df)
    fill_exam_slides(exam_slides, framing_df, exams_df)
    fill_legal_summary_slide(slide_legal, framing_df)
    fill_control_hierarchy_slide(slide_controls, actions)
    fill_action_slides(action_slides, actions)
    fill_timeline_slide(slide_timeline, actions, settings)
    fill_contact_slide(slide_contact, settings)

    company = safe_str(company_info.get("empresa"), "EMPRESA")
    period = safe_str(settings.get("period"), "PERÍODO")
    replace_common_placeholders(prs, company, period)

    # Se não houver logo, remove apenas o texto "INSERIR LOGO" para o arquivo final.
    if logo_bytes:
        for slide in prs.slides:
            add_logo_to_slide(slide, logo_bytes)
    else:
        for slide in prs.slides:
            logo_text = first_named(slide, "logo-text")
            if logo_text is not None:
                set_shape_text(logo_text, "")

    update_slide_numbers(prs)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# -----------------------------------------------------------------------------
# Interface Streamlit
# -----------------------------------------------------------------------------

st.set_page_config(page_title="Fortneer | Devolutiva SST", page_icon="📊", layout="wide")

st.title("Devolutiva Executiva de SST")
st.caption("Extração determinística do relatório + preenchimento manual do plano de ação + geração do PPTX Fortneer. Sem IA.")

if not DEFAULT_TEMPLATE.exists():
    st.error(
        "O arquivo Modelo_Devolutiva_SST_Fortneer.pptx não está na mesma pasta do app.py. "
        "Coloque o template ao lado do código e reinicie o Streamlit."
    )
    st.stop()

uploaded_report = st.file_uploader("1. Relatório Geral da empresa (PDF)", type=["pdf"])
logo_file = st.file_uploader("Logo do cliente (opcional)", type=["png", "jpg", "jpeg"])

if uploaded_report is None:
    st.info("Envie o Relatório de Empresas Geral para iniciar.")
    st.stop()

try:
    report_text = extract_pdf_text(uploaded_report.getvalue())
    company_auto, framing_auto, exams_auto = parse_report(report_text)
except Exception as exc:
    st.exception(exc)
    st.stop()

if framing_auto.empty:
    st.error("Nenhum bloco SETOR/CARGO foi identificado. Verifique se o PDF segue o mesmo padrão do relatório-modelo.")
    st.stop()

st.success(
    f"Relatório lido: {len(framing_auto)} cargos em {framing_auto['Setor'].nunique()} setores; "
    f"{len(exams_auto)} vínculos cargo/exame identificados."
)

tab_company, tab_framing, tab_exams, tab_actions, tab_output = st.tabs(
    ["Empresa", "Enquadramentos", "Exames", "Ações do PGR", "Gerar PPTX"]
)

with tab_company:
    c1, c2, c3 = st.columns(3)
    with c1:
        company_name = st.text_input("Empresa", value=company_auto.get("empresa", ""))
        cnpj = st.text_input("CNPJ", value=company_auto.get("cnpj", ""))
    with c2:
        employees = st.text_input("Nº de empregados", value=company_auto.get("num_empregados", ""))
        cnae = st.text_input("CNAE", value=company_auto.get("cnae", ""))
    with c3:
        risk_degree = st.text_input("Grau de risco", value=company_auto.get("grau_risco", ""))
        units = st.number_input("Unidades abrangidas", min_value=1, value=1, step=1)

    period = st.text_input("Período de referência", value=pd.Timestamp.today().strftime("%m/%Y"))
    documents = st.text_input("Documentos considerados", value="Relatório Geral / PGR / PCMSO / LTCAT / LIP")
    visits = st.text_input("Visitas e levantamentos", value="conforme documentação técnica")
    technical_responsibles = st.text_input("Responsáveis técnicos", value="conforme documentos vigentes")

    st.markdown("#### Dados executivos")
    c4, c5 = st.columns(2)
    with c4:
        documentation_status = st.text_input("Status da documentação", value="Relatório processado")
        exam_status = st.text_input("Status dos exames", value="Matriz extraída do relatório")
    with c5:
        risk_control_status = st.text_input("Status dos controles", value="Ver riscos e plano de ação")
        priority_impact = st.text_input(
            "Impacto esperado da prioridade",
            value="Impacto esperado: redução da exposição e regularização das pendências priorizadas.",
        )

with tab_framing:
    st.write(
        "O parser lê cada SETOR/CARGO e consolida as conclusões explícitas de insalubridade e os códigos GFIP encontrados nas CONCLUSÕES LTCAT. "
        "Revise a tabela antes de gerar."
    )
    framing_edit = st.data_editor(
        framing_auto[["Setor", "Cargo", "Insalubridade", "GFIP/LTCAT"]],
        key="framing_editor",
        use_container_width=True,
        hide_index=True,
        num_rows="fixed",
        column_config={
            "Setor": st.column_config.TextColumn(width="medium"),
            "Cargo": st.column_config.TextColumn(width="medium"),
            "Insalubridade": st.column_config.TextColumn(width="large"),
            "GFIP/LTCAT": st.column_config.TextColumn(width="medium"),
        },
    )

    with st.expander("Ver conclusões completas extraídas do relatório"):
        st.dataframe(
            framing_auto[["Setor", "Cargo", "Detalhe insalubridade", "Detalhe GFIP/LTCAT"]],
            use_container_width=True,
            hide_index=True,
        )

with tab_exams:
    st.write("ASO e Exames Obrigatórios são vinculados ao SETOR/CARGO e à periodicidade imediatamente associada no relatório.")
    exams_edit = st.data_editor(
        exams_auto[["Setor", "Cargo", "Exame", "Periodicidade"]],
        key="exams_editor",
        use_container_width=True,
        hide_index=True,
        num_rows="dynamic",
        column_config={
            "Setor": st.column_config.TextColumn(width="medium"),
            "Cargo": st.column_config.TextColumn(width="medium"),
            "Exame": st.column_config.TextColumn(width="medium"),
            "Periodicidade": st.column_config.TextColumn(width="large"),
        },
    )

with tab_actions:
    st.write(
        "Cadastre as ações do PGR. Os campos Grupo de risco e Setor/Cargo afetado alimentam o slide de riscos; "
        "Tipo de controle alimenta a hierarquia de controles; prazo alimenta o cronograma."
    )

    default_actions = pd.DataFrame(
        [
            {
                "Grupo de risco": "FÍSICO",
                "Risco/condição": "",
                "Setor/Cargo afetado": "",
                "Ação": "",
                "Prioridade": "ALTA",
                "Tipo de controle": "MEDIDA ADMINISTRATIVA",
                "Responsável": "",
                "Prazo (dias)": 30,
                "Status": "Aberta",
                "Evidência esperada": "",
            }
        ]
    )

    actions_edit = st.data_editor(
        default_actions,
        key="actions_editor",
        use_container_width=True,
        hide_index=True,
        num_rows="dynamic",
        column_config={
            "Grupo de risco": st.column_config.SelectboxColumn(options=RISK_GROUPS, required=False),
            "Risco/condição": st.column_config.TextColumn(width="medium"),
            "Setor/Cargo afetado": st.column_config.TextColumn(width="medium"),
            "Ação": st.column_config.TextColumn(width="large", required=False),
            "Prioridade": st.column_config.SelectboxColumn(options=["ALTA", "MÉDIA", "BAIXA"], required=True),
            "Tipo de controle": st.column_config.SelectboxColumn(options=CONTROL_TYPES, required=False),
            "Responsável": st.column_config.TextColumn(width="medium"),
            "Prazo (dias)": st.column_config.NumberColumn(min_value=0, step=1),
            "Status": st.column_config.SelectboxColumn(options=["Aberta", "Em curso", "Planejada", "Concluída"]),
            "Evidência esperada": st.column_config.TextColumn(width="large"),
        },
    )

with tab_output:
    st.markdown("#### Acompanhamento e contato")
    c1, c2 = st.columns(2)
    with c1:
        review_date = st.text_input("Data prevista para revisão", value="a definir")
        review_participants = st.text_input("Participantes da revisão", value="responsáveis definidos no plano")
    with c2:
        contact_name = st.text_input("Responsável Fortneer", value="")
        contact_phone = st.text_input("Telefone", value="")
        contact_email = st.text_input("E-mail", value="")

    normalized_actions = normalize_actions(actions_edit)

    st.markdown("#### Resumo antes de gerar")
    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Setores", framing_edit["Setor"].nunique())
    m2.metric("Cargos", len(framing_edit))
    m3.metric("Exames vinculados", len(exams_edit))
    m4.metric("Ações cadastradas", len(normalized_actions))

    positive_insal = framing_edit[
        ~framing_edit["Insalubridade"].isin(["Não insalubre", "Não identificado", ""])
    ]
    gfip_special = framing_edit[
        framing_edit["GFIP/LTCAT"].apply(lambda value: bool(extract_nontrivial_gfip(safe_str(value))))
    ]
    st.caption(
        f"Enquadramentos positivos de insalubridade: {len(positive_insal)} • "
        f"Cargos com GFIP diferente de 1: {len(gfip_special)}"
    )

    generate = st.button("Gerar apresentação PowerPoint", type="primary", use_container_width=True)

    if generate:
        company_info = {
            **company_auto,
            "empresa": company_name,
            "cnpj": cnpj,
            "num_empregados": employees,
            "cnae": cnae,
            "grau_risco": risk_degree,
        }
        settings = {
            "period": period,
            "units": units,
            "documents": documents,
            "visits": visits,
            "technical_responsibles": technical_responsibles,
            "documentation_status": documentation_status,
            "exam_status": exam_status,
            "risk_control_status": risk_control_status,
            "priority_impact": priority_impact,
            "review_date": review_date,
            "review_participants": review_participants,
            "contact_name": contact_name,
            "contact_phone": contact_phone,
            "contact_email": contact_email,
        }

        try:
            template_bytes = DEFAULT_TEMPLATE.read_bytes()
            pptx_bytes = build_presentation(
                template_bytes=template_bytes,
                company_info=company_info,
                framing_df=framing_edit.copy(),
                exams_df=exams_edit.copy(),
                actions_df=actions_edit.copy(),
                settings=settings,
                logo_bytes=logo_file.getvalue() if logo_file is not None else None,
            )

            safe_company = re.sub(r"[^A-Za-z0-9À-ÿ_-]+", "_", company_name).strip("_") or "Empresa"
            filename = f"Devolutiva_SST_{safe_company}.pptx"
            st.success("Apresentação gerada.")
            st.download_button(
                "Baixar PowerPoint",
                data=pptx_bytes,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        except Exception as exc:
            st.exception(exc)